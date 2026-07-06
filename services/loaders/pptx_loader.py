from pathlib import Path
from pptx import Presentation


def extract_pptx(file_path: Path):
    """
    Extract text from PowerPoint.
    """

    presentation = Presentation(file_path)

    slides_text = []

    for slide in presentation.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                text = shape.text.strip()

                if text:
                    slides_text.append(text)

    return {
        "filename": file_path.name,
        "document_type": "PowerPoint",
        "pages": len(presentation.slides),
        "text": "\n".join(slides_text)
    }


def extract_all_pptx(folder: Path):

    documents = []

    for file in folder.glob("*.pptx"):
        documents.append(extract_pptx(file))

    return documents